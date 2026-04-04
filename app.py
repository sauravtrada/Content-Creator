from flask import Flask, request, jsonify, render_template, send_file, url_for
from agent_graph import graph
import ppt_utils
import json
import os
import time
import tempfile
from rag_engine import ingest_document
from apscheduler.schedulers.background import BackgroundScheduler
import atexit

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", os.urandom(24))

# --- Background Cleanup Task ---
def cleanup_old_files():
    """Deletes .pptx files older than 1 hour."""
    now = time.time()
    cutoff = now - 3600  # 1 hour
    directory = os.path.dirname(os.path.abspath(__file__))
    for filename in os.listdir(directory):
        if filename.endswith(".pptx") and filename.startswith("presentation_"):
            filepath = os.path.join(directory, filename)
            try:
                if os.path.getmtime(filepath) < cutoff:
                    os.remove(filepath)
                    print(f"Deleted old file: {filename}")
            except Exception as e:
                print(f"Error checking/deleting {filename}: {e}")

# Initialize Scheduler
scheduler = BackgroundScheduler()
scheduler.add_job(func=cleanup_old_files, trigger="interval", minutes=60)
scheduler.start()

# Shut down the scheduler when exiting the app
atexit.register(lambda: scheduler.shutdown())

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate_ppt", methods=["POST"])
def generate_ppt():
    topic = request.form.get("topic")
    include_images = request.form.get("include_images", "false").lower() == "true"
    image_mode = request.form.get("image_mode", "manual")
    num_slides = int(request.form.get("num_slides", 5))
    tone = request.form.get("tone", "Professional")
    audience = request.form.get("audience", "General Audience")
    additional_instructions = request.form.get("additional_instructions", "")
    source_type = request.form.get("source_type", "topic_only")
    template_choice = request.form.get("template", "default")
    title_font_size = int(request.form.get("title_font_size", 26))
    body_font_size = int(request.form.get("body_font_size", 22))
    font_style = request.form.get("font_style", "Calibri")
    title_font_color = request.form.get("title_font_color", "#000000")
    body_font_color = request.form.get("body_font_color", "#333333")

    if not topic:
        return jsonify({"error": "Topic is required"}), 400

    try:
        rag_store = None
        
        if source_type == 'pdf':
            if 'file' not in request.files:
                return jsonify({"error": "No PDF file uploaded"}), 400
            file = request.files['file']
            if file.filename == '':
                return jsonify({"error": "No file selected"}), 400
            
            # Save temp file
            temp_dir = tempfile.gettempdir()
            temp_path = os.path.join(temp_dir, os.urandom(8).hex() + "_" + file.filename)
            file.save(temp_path)
            
            try:
                rag_store = ingest_document(source=temp_path, source_type='pdf')
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    
        elif source_type == 'url':
            source_url = request.form.get("source_url")
            if not source_url:
                return jsonify({"error": "URL is required"}), 400
            rag_store = ingest_document(source=source_url, source_type='url')
            
        elif source_type == 'text':
            source_text = request.form.get("source_text")
            if not source_text:
                return jsonify({"error": "Raw text is required"}), 400
            rag_store = ingest_document(source=source_text, source_type='text')

        # 1. Invoke LangGraph Workflow
        initial_state_and_prefs = {
            "initial_state": {
                "topic": topic,
                "source_type": source_type,
                "rag_store": rag_store,
                "include_images": include_images,
                "image_mode": image_mode,
                "num_slides": num_slides,
                "tone": tone,
                "audience": audience,
                "additional_instructions": additional_instructions,
                "outline": [],
                "slides": []
            },
            "design_prefs": {
                "template": template_choice,
                "title_font_size": title_font_size,
                "body_font_size": body_font_size,
                "font_style": font_style,
                "title_font_color": title_font_color,
                "body_font_color": body_font_color
            }
        }
        
        # 1. Invoke LangGraph Workflow
        result = graph.invoke(initial_state_and_prefs["initial_state"])
        json_content = result.get("final_output")
        
        if not json_content:
             raise ValueError("Graph failed to produce output")

        ppt_data = json.loads(json_content)
        
        # Inject design preferences into ppt_data for the util
        ppt_data["design_prefs"] = initial_state_and_prefs["design_prefs"]
        
        # 2. Create Presentation Locally
        filename = f"presentation_{os.urandom(4).hex()}.pptx"
        output_path = ppt_utils.create_ppt(
            ppt_data, 
            filename=filename, 
            image_mode=image_mode if include_images else None
        )

        # 3. Return the Download URL
        download_url = url_for('download_file', filename=filename)
        
        return jsonify({
            "message": "Presentation created successfully",
            "downloadUrl": download_url
        })

    except json.JSONDecodeError:
        return jsonify({"error": "Failed to generate valid JSON content from AI"}), 500
    except Exception as e:
        import traceback
        traceback.print_exc()
        error_msg = str(e)
        if "429" in error_msg or "RESOURCE_EXHAUSTED" in error_msg or "quota" in error_msg.lower():
            return jsonify({"error": "Gemini API Quota Exceeded. Please try again later or use a different API key/model."}), 429
        return jsonify({"error": error_msg}), 500


@app.route("/generate_content_only", methods=["POST"])
def generate_content_only():
    """Generates the JSON slide content without creating the PPTX file yet."""
    topic = request.form.get("topic")
    source_type = request.form.get("source_type", "topic_only")
    num_slides = int(request.form.get("num_slides", 5))
    include_images = request.form.get("include_images", "false").lower() == "true"
    image_mode = request.form.get("image_mode", "manual")
    tone = request.form.get("tone", "Professional")
    audience = request.form.get("audience", "General Audience")
    additional_instructions = request.form.get("additional_instructions", "")

    try:
        rag_store = None
        # Handle PDF/URL/Text ingestion (same as in generate_ppt)
        if source_type == 'pdf':
            if 'file' not in request.files:
                return jsonify({"error": "No PDF file uploaded"}), 400
            file = request.files['file']
            if file.filename == '':
                return jsonify({"error": "No file selected"}), 400
            
            temp_path = os.path.join(tempfile.gettempdir(), os.urandom(8).hex() + "_" + file.filename)
            file.save(temp_path)
            try: 
                rag_store = ingest_document(source=temp_path, source_type='pdf')
            finally: 
                if os.path.exists(temp_path): os.remove(temp_path)
        elif source_type == 'url':
            source_url = request.form.get("source_url")
            if not source_url:
                return jsonify({"error": "URL is required"}), 400
            rag_store = ingest_document(source=source_url, source_type='url')
        elif source_type == 'text':
            source_text = request.form.get("source_text")
            if not source_text:
                return jsonify({"error": "Raw text is required"}), 400
            rag_store = ingest_document(source=source_text, source_type='text')

        initial_state = {
            "topic": topic,
            "source_type": source_type,
            "rag_store": rag_store,
            "include_images": include_images,
            "image_mode": image_mode,
            "num_slides": num_slides,
            "tone": tone,
            "audience": audience,
            "additional_instructions": additional_instructions,
            "outline": [], "slides": []
        }
        
        result = graph.invoke(initial_state)
        json_content = result.get("final_output")
        if not json_content: raise ValueError("Graph failed")
        
        return json_content # This is already a JSON string from aggregator_node
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/update_slides", methods=["POST"])
def update_slides():
    """Takes existing slides and a user instruction, uses AI to update them."""
    try:
        from agent_graph import llm, extract_json
        
        current_slides_json = request.form.get("slides")
        instruction = request.form.get("instruction")
        
        if not current_slides_json or not instruction:
            return jsonify({"error": "Slides and instruction are required"}), 400
            
        current_slides = json.loads(current_slides_json)
        
        prompt = f"""
You are a professional presentation editor. 
Current Slides (JSON list):
{json.dumps(current_slides, indent=2)}

User Instruction:
"{instruction}"

Update the presentation according to the instruction. 
- You can MODIFY existing slides, DELETE slides, or **ADD NEW SLIDES** (by appending them to the list).
- If the user asks for a **NEW SLIDE**, create a new slide object with a logical heading and content.
- **IMAGES**: If the user asks for an image (on a new or existing slide):
  1. You MUST include: `"image_search_query": "a descriptive 3-5 word query for a high-quality photo"`
  2. You MUST set `"layout": "split"`.
- **VISUALS**: For any slide with a chart or table, set `"layout": "split"`. Otherwise use `"layout": "text_only"`.
- CHARTS: `"chart": {{"type": "column/pie/line/bar", "title": "Title", "categories": ["A", "B"], "series": [{{"name": "S1", "values": [10, 20]}}]}}`
- TABLES: `"table": [["Header 1", "Header 2"], ["Val 1", "Val 2"]]`
- STRUCTURE: Ensure the "content" field is a list of objects: `[{{"text": "point", "level": 0}}]`.
- Return ONLY the updated JSON list of slides. No preamble.

Updated JSON:
"""
        response = llm.invoke(prompt)
        updated_slides = extract_json(response.content)
        
        if not updated_slides:
            raise ValueError("Failed to parse updated slides from AI")
            
        # If AI wrapped it in an object like {"slides": [...]}, unwrap it
        if isinstance(updated_slides, dict) and "slides" in updated_slides:
            updated_slides = updated_slides["slides"]
        
        # --- NEW: Agentic Image Discovery for Updated Slides ---
        from agent_graph import image_searcher_node
        # Mock a state object for the node
        mock_state = {
            "slides": updated_slides,
            "include_images": True, 
            "image_mode": "manual" # Default to manual for updates
        }
        image_result = image_searcher_node(mock_state)
        updated_slides = image_result.get("slides", updated_slides)
            
        return jsonify(updated_slides)
    except Exception as e:
        import traceback
        traceback.print_exc()
        error_msg = str(e)
        if "429" in error_msg or "RESOURCE_EXHAUSTED" in error_msg or "quota" in error_msg.lower():
            return jsonify({"error": "Gemini API Quota Exceeded. Please try again later or use a different API key/model."}), 429
        return jsonify({"error": error_msg}), 500


@app.route("/generate_final", methods=["POST"])
def generate_final():
    """Takes (edited) JSON content and creates the final PPTX."""
    try:
        full_data_json = request.form.get("full_data")
        ppt_data = json.loads(full_data_json)
        
        # Inject design prefs
        ppt_data["design_prefs"] = {
            "template": request.form.get("template", "default"),
            "title_font_size": int(request.form.get("title_font_size", 28)),
            "body_font_size": int(request.form.get("body_font_size", 24)),
            "font_style": request.form.get("font_style", "Calibri"),
            "title_font_color": request.form.get("title_font_color", "#000000"),
            "body_font_color": request.form.get("body_font_color", "#333333")
        }

        filename = f"presentation_{os.urandom(4).hex()}.pptx"
        output_path = ppt_utils.create_ppt(
            ppt_data, 
            filename=filename, 
            image_mode=request.form.get("image_mode", "manual")
        )
        download_url = url_for('download_file', filename=filename)
        
        return jsonify({"message": "Success", "downloadUrl": download_url})
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/import_ppt", methods=["POST"])
def import_ppt():
    """Parses an uploaded PPTX and returns JSON slides."""
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    temp_path = os.path.join(tempfile.gettempdir(), os.urandom(8).hex() + "_" + file.filename)
    file.save(temp_path)
    
    try:
        from pptx import Presentation
        prs = Presentation(temp_path)
        slides = []
        for slide in prs.slides:
            heading = slide.shapes.title.text if slide.shapes.title else "Untitled Slide"
            content = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        content.append({"text": paragraph.text, "level": paragraph.level})
            slides.append({"heading": heading, "content": content})
            
        return jsonify({
            "title": file.filename.replace(".pptx", ""),
            "theme": {"font": "Calibri"},
            "slides": slides
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        if os.path.exists(temp_path): os.remove(temp_path)

@app.route('/download/<filename>')
def download_file(filename):
    # Determine the path. 
    # Current implementation of create_ppt saves to current working directory or absolute path.
    # ppt_utils.create_ppt saves to os.path.abspath(filename) if we passed just filename.
    # So valid check if file exists in current dir.
    # SECURITY NOTE: In production, sanitize filename to prevent directory traversal.
    file_path = os.path.abspath(filename)
    if os.path.exists(file_path):
        return send_file(
            file_path, 
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        return "File not found", 404

if __name__ == "__main__":
    app.run(debug=False)
