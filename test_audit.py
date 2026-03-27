import os
import json
from agent_graph import graph
import ppt_utils
from dotenv import load_dotenv

load_dotenv()

def test_flow():
    print("starting end-to-end test...")
    
    # Simulate the data passed from app.py
    topic = "The Future of Quantum Computing"
    initial_state = {
        "topic": topic,
        "source_type": "topic_only",
        "rag_store": None,
        "include_images": False,
        "image_mode": None,
        "num_slides": 3,
        "tone": "Professional",
        "audience": "Tech Enthusiasts",
        "additional_instructions": "Focus on practical applications.",
        "outline": [],
        "slides": []
    }
    
    design_prefs = {
        "template": "professional_blue.pptx",
        "title_font_size": 40,
        "body_font_size": 24,
        "title_font_color": "#000000",
        "body_font_color": "#333333"
    }

    print(f"invoking langgraph for topic: {topic}")
    try:
        # Run the workflow
        result = graph.invoke(initial_state)
        json_content = result.get("final_output")
        
        if not json_content:
            print("Error: No final output from graph")
            return

        print("Writing RAW JSON to raw_ai_output.json...")
        with open("raw_ai_output.json", "w") as f:
            f.write(json_content)
            
        ppt_data = json.loads(json_content)
        print(f"PARSED PPT DATA: title='{ppt_data.get('title')}', slides_count={len(ppt_data.get('slides', []))}")
        ppt_data["design_prefs"] = design_prefs
        
        print("generating .pptx file...")
        filename = "test_audit_output.pptx"
        output_path = ppt_utils.create_ppt(
            ppt_data, 
            filename=filename, 
            image_mode=None
        )
        
        print(f"success! presentation saved to: {output_path}")
        
        # Basic validation of the file
        if os.path.exists(output_path):
            from pptx import Presentation
            prs = Presentation(output_path)
            print(f"validated .pptx: found {len(prs.slides)} slides.")
            for i, slide in enumerate(prs.slides):
                print(f"  slide {i}: layout='{slide.slide_layout.name}'")
        
    except Exception as e:
        print(f"test failed with error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_flow()
