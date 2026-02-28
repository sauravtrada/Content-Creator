import json
import os
from pptx import Presentation

def analyze_ppt(filepath):
    prs = Presentation(filepath)
    results = {
        "slide_size": {"width": prs.slide_width.inches, "height": prs.slide_height.inches},
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_info = {"slide_number": i + 1, "shapes": [], "warnings": []}
        shapes_info = []
        
        for j, shape in enumerate(slide.shapes):
            shape_data = {
                "id": j,
                "left": shape.left.inches if shape.left else 0,
                "top": shape.top.inches if shape.top else 0,
                "width": shape.width.inches if shape.width else 0,
                "height": shape.height.inches if shape.height else 0,
            }
            shape_data["right"] = shape_data["left"] + shape_data["width"]
            shape_data["bottom"] = shape_data["top"] + shape_data["height"]
            
            if shape.has_text_frame:
                shape_data["type"] = "text"
                shape_data["text"] = shape.text.replace("\n", " ").strip()[:100]
            elif shape.shape_type == 13: # Picture
                shape_data["type"] = "picture"
            else:
                shape_data["type"] = f"other_{shape.shape_type}"
                
            slide_info["shapes"].append(shape_data)
            shapes_info.append(shape_data)
            
            # Check for out of bounds
            if shape_data["left"] < 0 or shape_data["top"] < 0 or shape_data["right"] > prs.slide_width.inches or shape_data["bottom"] > prs.slide_height.inches:
                slide_info["warnings"].append(f"Shape {j} ({shape_data['type']}) is OUT OF BOUNDS.")
                
        # Check for overlaps
        for a in range(len(shapes_info)):
            for b in range(a + 1, len(shapes_info)):
                s1 = shapes_info[a]
                s2 = shapes_info[b]
                # Bounding box overlap check
                if not (s1['right'] <= s2['left'] or s1['left'] >= s2['right'] or s1['bottom'] <= s2['top'] or s1['top'] >= s2['bottom']):
                    # Ignore background shapes or full slide shapes
                    if s1['width'] > 10 and s1['height'] > 6: continue
                    if s2['width'] > 10 and s2['height'] > 6: continue
                    slide_info["warnings"].append(f"Shape {a} ('{s1.get('text', s1['type'])}') and Shape {b} ('{s2.get('text', s2['type'])}') OVERLAP.")
                    
        results["slides"].append(slide_info)

    with open("ppt_analysis.json", "w", encoding="utf-8") as f:
        json.dump(results, f, indent=2)

if __name__ == "__main__":
    analyze_ppt("test_layout_theme.pptx")
