from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import requests
from io import BytesIO
import concurrent.futures
import urllib.parse
from datetime import datetime
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import random

# --- Layout Constants ---
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT  = Inches(0.9)
MARGIN_RIGHT = Inches(0.9)
MARGIN_TOP   = Inches(0.1)


# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str, default=None):
    """Convert a hex colour string to RGBColor. Returns *default* on failure."""
    if not hex_str:
        return default
    try:
        h = str(hex_str).replace("#", "").strip()
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) == 6:
            return RGBColor.from_string(h.upper())
    except Exception:
        pass
    return default


def fetch_image(url_or_query):
    """
    Fetches an image stream using the tiered image service.
    Delegates to services/image_service.py:
      Unsplash → Pexels → DuckDuckGo → None
    Returns BytesIO or None. Never raises.
    """
    if not url_or_query:
        return None
    try:
        from services.image_service import fetch_image_for_query
        print(f"--- DEBUG: Fetching image for: {url_or_query[:60]}")
        return fetch_image_for_query(url_or_query)
    except Exception as e:
        print(f"--- DEBUG: fetch_image error: {e}")
        return None


# ---------------------------------------------------------------------------
# Theme registry — matches the 12 templates in slide_templates/
# ---------------------------------------------------------------------------
THEME_REGISTRY = {
    "professional_blue.pptx": {
        "bg": "#FFFFFF", "surface": "#EBF3FF", "accent": "#1565C0", "text": "#212121",
    },
    "executive_dark.pptx": {
        "bg": "#1A1A2E", "surface": "#16213E", "accent": "#E94560", "text": "#EAEAEA",
    },
    "teal_modern.pptx": {
        "bg": "#F0FDFA", "surface": "#CCFBF1", "accent": "#0D9488", "text": "#134E4A",
    },
    "sunset_coral.pptx": {
        "bg": "#FFF5F5", "surface": "#FED7D7", "accent": "#E53E3E", "text": "#2D3748",
    },
    "ocean_blue.pptx": {
        "bg": "#EFF6FF", "surface": "#DBEAFE", "accent": "#2563EB", "text": "#1E3A8A",
    },
    "forest_green.pptx": {
        "bg": "#F0FDF4", "surface": "#DCFCE7", "accent": "#15803D", "text": "#14532D",
    },
    "royal_purple.pptx": {
        "bg": "#1E1B4B", "surface": "#312E81", "accent": "#818CF8", "text": "#EDE9FE",
    },
    "golden_amber.pptx": {
        "bg": "#FFFBEB", "surface": "#FEF3C7", "accent": "#D97706", "text": "#451A03",
    },
    "charcoal_mono.pptx": {
        "bg": "#18181B", "surface": "#27272A", "accent": "#A1A1AA", "text": "#FAFAFA",
    },
    "rose_blush.pptx": {
        "bg": "#FFF1F2", "surface": "#FFE4E6", "accent": "#E11D48", "text": "#1F2937",
    },
    "navy_slate.pptx": {
        "bg": "#0F172A", "surface": "#1E293B", "accent": "#38BDF8", "text": "#F0F9FF",
    },
    "olive_sage.pptx": {
        "bg": "#FAFAF9", "surface": "#F5F5F4", "accent": "#65A30D", "text": "#1A2E05",
    },
    "nebula_glow.pptx": {
        "bg": "#0D0D15", "surface": "#1A1A2E", "accent": "#7AA2F7", "text": "#E0E0FB",
    },
    "retrowave_neon.pptx": {
        "bg": "#240046", "surface": "#3C096C", "accent": "#FF00FF", "text": "#FFFFFF",
    },
    "facet.pptx": {
        "bg": "#FFFFFF", "surface": "#F0F8F8", "accent": "#008B8B", "text": "#2F4F4F",
    },
    "gallery.pptx": {
        "bg": "#FFFFFF", "surface": "#FDF5E6", "accent": "#A52A2A", "text": "#3E2723",
    },
    "integral.pptx": {
        "bg": "#FFFFFF", "surface": "#F5FFFA", "accent": "#2E8B57", "text": "#1B5E20",
    },
    "ion.pptx": {
        "bg": "#FFFFFF", "surface": "#F8F8FF", "accent": "#6A5ACD", "text": "#1A237E",
    },
    "berlin.pptx": {
        "bg": "#FFFFFF", "surface": "#FFF8DC", "accent": "#D2691E", "text": "#3E2723",
    },
}


def _solid_rect(slide, left, top, width, height, color):
    """Add a solid filled, borderless rectangle to *slide*."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _push_back(slide, element, index=2):
    """Move *element* behind existing shapes at z-order *index*."""
    slide.shapes._spTree.remove(element)
    slide.shapes._spTree.insert(index, element)


def _add_chart(slide, chart_data_dict, x, y, w, h):
    """Adds a chart to the slide based on data dictionary."""
    try:
        chart_type_str = chart_data_dict.get("type", "column").lower()
        title = chart_data_dict.get("title", "")
        categories = chart_data_dict.get("categories", [])
        series_list = chart_data_dict.get("series", [])

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series_list:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type_str == "pie":
            chart_type = XL_CHART_TYPE.PIE
        elif chart_type_str == "line":
            chart_type = XL_CHART_TYPE.LINE
        elif chart_type_str == "bar":
            chart_type = XL_CHART_TYPE.BAR_CLUSTERED

        chart_shape = slide.shapes.add_chart(
            chart_type, x, y, w, h, chart_data
        )
        if title:
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.text_frame.text = title
        return chart_shape
    except Exception as e:
        print(f"Error adding chart: {e}")
        return None


def _add_table(slide, table_data_list, x, y, w, h):
    """Adds a table to the slide based on list of lists."""
    try:
        rows = len(table_data_list)
        cols = len(table_data_list[0]) if rows > 0 else 0
        if rows == 0 or cols == 0:
            return None

        table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
        table = table_shape.table

        for r_idx, row_data in enumerate(table_data_list):
            for c_idx, cell_val in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_val)
                # Apply consistent styling to all cells
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)  # Smaller font to fit more rows
                    paragraph.font.name = "Calibri"
                    if r_idx == 0:
                        paragraph.font.bold = True
        return table_shape
    except Exception as e:
        print(f"Error adding table: {e}")
        return None


def set_slide_background(slide, color):
    """Explicitly set the slide background color property."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def clear_all_slides(prs):
    """Remove all slides from the presentation object."""
    # iterate backwards to avoid index issues
    xml_slides = prs.slides._sldIdLst
    for i in range(len(xml_slides) - 1, -1, -1):
        xml_slides.remove(xml_slides[i])

def get_layout_by_name(prs, name_substring):
    """Finds a slide layout where the name contains the substring (case-insensitive)."""
    for layout in prs.slide_layouts:
        if name_substring.lower() in layout.name.lower():
            return layout
    # Fallback to a common index if name is not found, but log it
    print(f"Warning: Layout with name '{name_substring}' not found. Falling back to index-based selection.")
    if "title" in name_substring.lower() and "content" in name_substring.lower():
        return prs.slide_layouts[1]
    if "two" in name_substring.lower() or "comparison" in name_substring.lower():
        return prs.slide_layouts[3]
    return prs.slide_layouts[0]


# ---------------------------------------------------------------------------
# Pagination helper
# ---------------------------------------------------------------------------

def _paginate_slides(original_slides, image_mode, user_body_size):
    MAX_HEIGHT_PT    = 480  # Increased for more content per slide
    ITEM_SPACING     = {0: 12, 1: 8, 2: 4} # Reduced spacing to fit more items
    FONT_MAP         = {0: user_body_size, 1: max(12, user_body_size - 4), 2: max(10, user_body_size - 6)}
    
    def get_chars_limit(level, has_visual):
        font_size = FONT_MAP.get(level, 24)
        # Full area ~900pt, Split area ~450pt. Heuristic char width ~0.55*font
        area_width = 450 if has_visual else 900
        return max(20, int(area_width / (font_size * 0.55)))

    def estimate_h(text, level, has_visual):
        if not text:
            return 0
        chars_limit = get_chars_limit(level, has_visual)
        lines   = max(1, (len(text) + chars_limit - 1) // chars_limit)
        line_h  = FONT_MAP.get(level, 24) * 1.2
        return lines * line_h + ITEM_SPACING.get(level, 10)

    paginated = []
    for slide_data in original_slides:
        # Determine if this specific slide will have a second column (visual)
        explicit_layout = slide_data.get("layout", "auto").lower()
        if explicit_layout == "text_only":
            has_visual = False
        elif explicit_layout == "split":
            has_visual = True
        else:
            # Fallback: Split if chart, table, or an image query exists.
            has_visual = (
                bool(slide_data.get("chart")) or 
                bool(slide_data.get("table")) or 
                bool(slide_data.get("image_search_query"))
            )

        content = slide_data.get("content", [])
        if not content and "bullet_points" in slide_data:
            content = [{"text": bp, "level": 0} for bp in slide_data["bullet_points"]]
        
        # Estimate table height if present
        table_data = slide_data.get("table", [])
        table_h = len(table_data) * 28 if table_data else 0 # ~28pt per row with 12pt font

        if not content:
            paginated.append(slide_data)
            continue

        current_items  = []
        current_height = table_h # Start with table height
        for item in content:
            h = estimate_h(item.get("text", ""), item.get("level", 0), has_visual)
            if current_height + h > MAX_HEIGHT_PT and current_items:
                chunk = slide_data.copy()
                chunk["content"] = current_items
                paginated.append(chunk)
                current_items  = [item]
                current_height = h # New slide (Cont.) doesn't have the table
            else:
                current_items.append(item)
                current_height += h

        if current_items:
            chunk = slide_data.copy()
            chunk["content"] = current_items
            base_title = slide_data.get("heading", "")
            if paginated and paginated[-1].get("heading", "").replace(" (Cont.)", "") == base_title.replace(" (Cont.)", ""):
                chunk["heading"] = f"{base_title} (Cont.)"
                chunk.pop("image_search_query", None)
                chunk.pop("chart", None)
                chunk.pop("table", None)
            paginated.append(chunk)

    return paginated


# ---------------------------------------------------------------------------
# Main PPT creation function
# ---------------------------------------------------------------------------

def create_ppt(data, filename="presentation.pptx", image_mode="manual"):
    """
    The main driver: takes the JSON data and produces a .pptx file.
    """
    print(f"--- DEBUG: ppt_utils.create_ppt called for {filename} ---")
    design_prefs   = data.get("design_prefs", {})
    template_name  = design_prefs.get("template", "professional_blue.pptx")
    template_path  = os.path.join("slide_templates", template_name)

    # --- Load or Create Presentation ---
    template_path = os.path.join("slide_templates", template_name)
    if not os.path.exists(template_path):
        template_path = os.path.join("slide_templates", "professional_blue.pptx")

    try:
        prs = Presentation(template_path)
        # CRITICAL: Clear any existing slides from the template to prevent relationship collisions
        clear_all_slides(prs)
    except Exception as e:
        print(f"Error loading template {template_path}: {e}")
        prs = Presentation()
    
    # Always ensure slide dimensions are correct
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    W = prs.slide_width
    H = prs.slide_height

    presentation_title = data.get("title", "Untitled Presentation")

    # --- Dynamic font sizes ---
    user_title_size = design_prefs.get("title_font_size", 40)
    user_body_size  = design_prefs.get("body_font_size", 24)

    DYN_TITLE_SIZE   = Pt(user_title_size + 14)
    DYN_SUBTITLE_SIZE = Pt(max(16, user_title_size - 10))
    DYN_HEADING_SIZE  = Pt(user_title_size)
    DYN_BODY_L0       = Pt(user_body_size)
    DYN_BODY_L1       = Pt(max(12, user_body_size - 4))
    DYN_BODY_L2       = Pt(max(10, user_body_size - 6))

    # --- Resolve colours (theme registry → user overrides) ---
    registry    = THEME_REGISTRY.get(template_name, THEME_REGISTRY["professional_blue.pptx"])
    THEME_FONT  = design_prefs.get("font_style", data.get("theme", {}).get("font", "Calibri"))

    C_BG      = hex_to_rgb(registry["bg"],      RGBColor(255, 255, 255))
    C_SURFACE = hex_to_rgb(registry["surface"],  RGBColor(235, 243, 255))
    C_ACCENT  = hex_to_rgb(registry["accent"],   RGBColor(21,  101, 192))
    C_TEXT    = hex_to_rgb(registry["text"],     RGBColor(33,  33,  33))

    # Allow per-field user overrides from the form
    C_TITLE   = hex_to_rgb(design_prefs.get("title_font_color"),
                             C_ACCENT)  # default to accent
    C_BODY    = hex_to_rgb(design_prefs.get("body_font_color"), C_TEXT)
    C_FOOTER  = RGBColor(
        min(255, C_TEXT[0] + 60),
        min(255, C_TEXT[1] + 60),
        min(255, C_TEXT[2] + 60),
    )

    # -----------------------------------------------------------------------
    # apply_title_slide_styling  — split-panel + accents
    # -----------------------------------------------------------------------
    def apply_title_slide_styling(slide):
        set_slide_background(slide, C_BG)
        bg_shp = _solid_rect(slide, 0, 0, W, H, C_BG)
        _push_back(slide, bg_shp._element, 2)

        if template_name == "nebula_glow.pptx":
            # Central floating card (glassish)
            card_w, card_h = int(W * 0.7), int(H * 0.6)
            card = _solid_rect(slide, (W-card_w)//2, (H-card_h)//2, card_w, card_h, C_SURFACE)
            try:
                card.fill.fore_color.brightness = 0.1
            except: pass
            _push_back(slide, card._element, 3)
            # Corner accents
            sz = Inches(1.5)
            _push_back(slide, _solid_rect(slide, 0, 0, sz, sz, C_ACCENT)._element, 4)
            _push_back(slide, _solid_rect(slide, W-sz, H-sz, sz, sz, C_ACCENT)._element, 4)
        
        elif template_name == "retrowave_neon.pptx":
            # Grid lines at bottom
            grid_h = int(H * 0.3)
            grid = _solid_rect(slide, 0, H-grid_h, W, grid_h, C_SURFACE)
            _push_back(slide, grid._element, 3)
            # Accent line (horizon)
            _push_back(slide, _solid_rect(slide, 0, H-grid_h-Inches(0.05), W, Inches(0.05), C_ACCENT)._element, 4)
            # Decorative sun (square proxy for now as add_shape(1) is rect)
            sun_sz = Inches(3)
            sun = _solid_rect(slide, (W-sun_sz)//2, Inches(1), sun_sz, sun_sz, C_ACCENT)
            _push_back(slide, sun._element, 3)

        elif template_name == "facet.pptx":
            # Triangular shape in top left and bottom right
            tri_sz = Inches(2)
            _push_back(slide, _solid_rect(slide, 0, 0, tri_sz, tri_sz, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, W-tri_sz, H-tri_sz, tri_sz, tri_sz, C_ACCENT)._element, 3)
            # Vertical line
            _push_back(slide, _solid_rect(slide, Inches(0.5), 0, Inches(0.05), H, C_SURFACE)._element, 4)

        elif template_name == "gallery.pptx":
            # Canvas-like border
            border_w = Inches(0.4)
            _push_back(slide, _solid_rect(slide, 0, 0, W, border_w, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, 0, H-border_w, W, border_w, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, 0, 0, border_w, H, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, W-border_w, 0, border_w, H, C_ACCENT)._element, 3)

        elif template_name == "integral.pptx":
            # Strong vertical band on the left
            band_w = Inches(1.2)
            _push_back(slide, _solid_rect(slide, 0, 0, band_w, H, C_ACCENT)._element, 3)
            # Thin horizontal line near top
            _push_back(slide, _solid_rect(slide, 0, Inches(1.5), W, Inches(0.05), C_SURFACE)._element, 4)

        elif template_name == "ion.pptx":
            # Minimalist horizontal bars
            bar_h = Inches(0.3)
            _push_back(slide, _solid_rect(slide, 0, H-bar_h, W, bar_h, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, Inches(1), Inches(1), Inches(4), Inches(0.1), C_ACCENT)._element, 4)

        elif template_name == "berlin.pptx":
            # Solid block header and footer
            header_h = Inches(1)
            _push_back(slide, _solid_rect(slide, 0, 0, W, header_h, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, 0, H-Inches(0.5), W, Inches(0.5), C_SURFACE)._element, 4)
            # Accent square
            _push_back(slide, _solid_rect(slide, Inches(0.5), Inches(0.5), Inches(1), Inches(1), C_BG)._element, 5)

        else:
            # Default logic (Original)
            band_w = int(W * 0.38)
            band   = _solid_rect(slide, 0, 0, band_w, H, C_ACCENT)
            _push_back(slide, band._element, 3)
            corner_sz = Inches(1.8)
            corner    = _solid_rect(slide, W - corner_sz, 0, corner_sz, corner_sz, C_ACCENT)
            try:
                corner.fill.fore_color.brightness = 0.25
            except: pass
            _push_back(slide, corner._element, 3)
            edge_h = Inches(0.05)
            edge   = _solid_rect(slide, 0, H - edge_h, W, edge_h, C_ACCENT)
            _push_back(slide, edge._element, 3)

    def apply_content_slide_styling(slide):
        set_slide_background(slide, C_BG)
        bg_shp = _solid_rect(slide, 0, 0, W, H, C_BG)
        _push_back(slide, bg_shp._element, 2)

        if template_name == "nebula_glow.pptx":
            # Top-right glowing blob
            blob_sz = Inches(2.5)
            blob = _solid_rect(slide, W - blob_sz, 0, blob_sz, blob_sz, C_ACCENT)
            _push_back(slide, blob._element, 3)
            # Thin bottom accent
            _push_back(slide, _solid_rect(slide, Inches(0.5), H-Inches(0.4), W-Inches(1), Inches(0.03), C_ACCENT)._element, 3)
        
        elif template_name == "retrowave_neon.pptx":
            # Sidebar accent
            _push_back(slide, _solid_rect(slide, 0, 0, Inches(0.2), H, C_ACCENT)._element, 3)
            # Header block
            _push_back(slide, _solid_rect(slide, Inches(0.2), 0, W-Inches(0.2), Inches(1.2), C_SURFACE)._element, 3)
        
        elif template_name == "facet.pptx":
            # Top accent bar
            _push_back(slide, _solid_rect(slide, 0, 0, W, Inches(0.15), C_ACCENT)._element, 3)
            # Subtle left indicator
            _push_back(slide, _solid_rect(slide, 0, 0, Inches(0.08), H, C_SURFACE)._element, 3)

        elif template_name == "gallery.pptx":
            # Minimal thin frame
            frame_sz = Inches(0.05)
            _push_back(slide, _solid_rect(slide, 0, 0, W, frame_sz, C_ACCENT)._element, 3)
            _push_back(slide, _solid_rect(slide, 0, H-frame_sz, W, frame_sz, C_ACCENT)._element, 3)

        elif template_name == "integral.pptx":
            # Left vertical marker
            _push_back(slide, _solid_rect(slide, 0, 0, Inches(0.6), H, C_ACCENT)._element, 3)

        elif template_name == "ion.pptx":
            # Bottom accent line
            _push_back(slide, _solid_rect(slide, 0, H-Inches(0.12), W, Inches(0.12), C_ACCENT)._element, 3)

        elif template_name == "berlin.pptx":
            # Header strip
            _push_back(slide, _solid_rect(slide, 0, 0, W, Inches(1.15), C_SURFACE)._element, 3)
            _push_back(slide, _solid_rect(slide, 0, Inches(1.15), W, Inches(0.04), C_ACCENT)._element, 4)

        else:
            # Default logic (Original)
            bar_w = Inches(0.12)
            bar   = _solid_rect(slide, 0, 0, bar_w, H, C_ACCENT)
            _push_back(slide, bar._element, 3)
            strip_h = Inches(1.1)
            strip   = _solid_rect(slide, bar_w, 0, W - bar_w, strip_h, C_SURFACE)
            _push_back(slide, strip._element, 3)
            uline_h = Inches(0.045)
            uline   = _solid_rect(slide, bar_w, strip_h - uline_h, W - bar_w, uline_h, C_ACCENT)
            _push_back(slide, uline._element, 4)
            footer_y = H - Inches(0.52)
            footer_line = _solid_rect(slide, bar_w + Inches(0.3), footer_y, W - bar_w - Inches(0.6), Inches(0.022), C_ACCENT)
            _push_back(slide, footer_line._element, 4)

    # -----------------------------------------------------------------------
    # Footer text
    # -----------------------------------------------------------------------
    def add_footer(slide, slide_number):
        date_str = datetime.now().strftime("%B %d, %Y")
        # Slide number — bottom right
        sn_w = Inches(1)
        sn_h = Inches(0.4)
        sn_l = W - sn_w - Inches(0.25)
        sn_t = H - sn_h - Inches(0.07)
        tb   = slide.shapes.add_textbox(sn_l, sn_t, sn_w, sn_h)
        p    = tb.text_frame.paragraphs[0]
        p.text              = str(slide_number)
        p.font.size         = Pt(11)
        p.font.color.rgb    = C_FOOTER
        p.alignment         = PP_ALIGN.RIGHT

        # Footer label — bottom left
        fl_w = Inches(9)
        fl_l = Inches(0.5)
        tb2  = slide.shapes.add_textbox(fl_l, sn_t, fl_w, sn_h)
        p2   = tb2.text_frame.paragraphs[0]
        p2.text              = f"{presentation_title}  |  {date_str}"
        p2.font.size         = Pt(11)
        p2.font.color.rgb    = C_FOOTER
        p2.alignment         = PP_ALIGN.LEFT

    # -----------------------------------------------------------------------
    # SLIDE 1 — Cover / Title Slide
    # -----------------------------------------------------------------------
    title_layout = get_layout_by_name(prs, "Title Slide")

    cover = prs.slides.add_slide(title_layout)
    apply_title_slide_styling(cover)

    # Title text (right panel — place over the white/bg area)
    band_w_val = int(W * 0.38)
    right_x    = band_w_val + Inches(0.5)
    right_w    = W - band_w_val - Inches(1)

    title_ph = cover.shapes.title
    title_ph.left   = right_x
    title_ph.top    = int(H * 0.28)
    title_ph.width  = right_w
    title_ph.height = Inches(2.2)
    title_ph.text   = presentation_title
    tf = title_ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.bold      = True
    p.font.size      = DYN_TITLE_SIZE
    p.font.name      = THEME_FONT
    p.font.color.rgb = C_TITLE
    p.alignment      = PP_ALIGN.LEFT

    # Subtitle
    try:
        sub_ph = cover.placeholders[1]
        sub_ph.left   = right_x
        sub_ph.top    = int(H * 0.28) + Inches(2.3)
        sub_ph.width  = right_w
        sub_ph.height = Inches(0.8)
        sub_ph.text   = "Generated by Gemini AI"
        sp = sub_ph.text_frame.paragraphs[0]
        sp.font.size         = DYN_SUBTITLE_SIZE
        sp.font.name         = THEME_FONT
        sp.font.color.rgb    = C_BODY
        sp.alignment         = PP_ALIGN.LEFT
    except Exception:
        pass

    # -----------------------------------------------------------------------
    # Content Slides
    # -----------------------------------------------------------------------
    final_slides = _paginate_slides(
        data.get("slides", []), image_mode, user_body_size
    )

    # Pre-fetch images in parallel if auto mode or if explicit queries are present
    image_map = {}
    if image_mode:
        print(f"Fetching images for {len(final_slides)} slides...")
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            future_to_i = {}
            for i, sd in enumerate(final_slides):
                # Priority 0: Explicit image URL found by agent
                # Priority 1: Explicit agent-provided query
                # Priority 2: In 'auto' mode, use slide heading as fallback
                url = sd.get("image_url")
                query = sd.get("image_search_query")
                if not query and image_mode == "auto":
                    query = sd.get("heading", "")
                
                target = url if url else query
                if target:
                    future_to_i[executor.submit(fetch_image, target)] = i
            
            for future in concurrent.futures.as_completed(future_to_i):
                idx = future_to_i[future]
                try:
                    stream = future.result()
                    if stream:
                        image_map[idx] = stream
                except Exception as exc:
                    print(f"Image error slide {idx}: {exc}")

    # Choose layout
    use_image = image_mode in ("manual", "auto")

    for i, slide_data in enumerate(final_slides):
        # Explicit layout support
        explicit_layout = slide_data.get("layout", "auto").lower()
        
        if explicit_layout == "text_only":
            has_any_visual = False
        elif explicit_layout == "split":
            has_any_visual = True
        else:
            # Fallback to automatic detection
            # Split if a chart, table, or an image query exists
            has_chart = bool(slide_data.get("chart"))
            has_table = bool(slide_data.get("table"))
            has_image_query = bool(slide_data.get("image_search_query"))
            has_auto_image = (image_mode == "auto" and i in image_map)
            
            has_any_visual = has_chart or has_table or has_image_query or has_auto_image

        if has_any_visual:
            layout = get_layout_by_name(prs, "Two Content")
        else:
            layout = get_layout_by_name(prs, "Title and Content")
            
        slide  = prs.slides.add_slide(layout)
        apply_content_slide_styling(slide)

        shapes = slide.shapes
        BAR_W  = Inches(0.12)
        STRIP_H = Inches(1.1)

        # --- Title ---
        title_shp        = shapes.title
        title_shp.left   = BAR_W + Inches(0.3)
        title_shp.top    = Inches(0.12)
        title_shp.height = Inches(0.85)
        
        # Title always uses full width to avoid unnecessary empty space on the right
        title_shp.width = W - BAR_W - Inches(0.7)

        title_shp.text = slide_data.get("heading", "Slide")
        tf = title_shp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.bold      = True
        p.font.size      = DYN_HEADING_SIZE
        p.font.name      = THEME_FONT
        p.font.color.rgb = C_TITLE
        p.alignment      = PP_ALIGN.LEFT

        # --- Footer ---
        add_footer(slide, i + 1)

        # --- Body ---
        body_shp        = shapes.placeholders[1]
        body_shp.left   = BAR_W + Inches(0.3)
        body_shp.top    = STRIP_H + Inches(0.15)

        if has_any_visual:
            # Leave 50% for visuals on the right
            body_shp.width = int((W - BAR_W - Inches(0.7)) * 0.50)
        else:
            body_shp.width = W - BAR_W - Inches(0.7)

        body_shp.height = H - body_shp.top - Inches(0.6)

        tf2 = body_shp.text_frame
        tf2.clear()
        tf2.word_wrap = True

        content = slide_data.get("content", [])
        if not content and "bullet_points" in slide_data:
            content = [{"text": bp, "level": 0} for bp in slide_data["bullet_points"]]

        for item in content:
            text  = item.get("text", "")
            level = item.get("level", 0)

            para = tf2.add_paragraph()
            para.text  = text
            para.level = level
            para.font.name = THEME_FONT
            para.alignment = PP_ALIGN.LEFT

            if level == 0:
                para.font.size      = DYN_BODY_L0
                para.font.bold      = False
                para.font.color.rgb = C_BODY
                para.space_before   = Pt(6)
                para.space_after    = Pt(4)
            elif level == 1:
                para.font.size      = DYN_BODY_L1
                para.font.bold      = False
                para.font.color.rgb = C_BODY
                para.space_before   = Pt(4)
                para.space_after    = Pt(2)
            else:
                para.font.size      = DYN_BODY_L2
                para.font.bold      = False
                para.font.color.rgb = C_BODY
                para.space_before   = Pt(2)
                para.space_after    = Pt(2)

        # --- Image ---
        visual_x = body_shp.left + body_shp.width + Inches(0.4)
        visual_y = body_shp.top
        visual_w = W - visual_x - Inches(0.4)
        visual_h = body_shp.height

        # Fix: Previously only inserted if mode == 'auto'. 
        # Now inserts if there's an image in map (which handles both auto and manual queries).
        if use_image and i in image_map:
            try:
                stream = image_map[i]
                stream.seek(0)
                pic = slide.shapes.add_picture(stream, visual_x, visual_y, width=visual_w, height=visual_h)
                
                # To ensure it fills exactly, we set the position and size again
                pic.left   = visual_x
                pic.top    = visual_y
                pic.width  = visual_w
                pic.height = visual_h
            except Exception as e:
                import traceback
                print(f"--- DEBUG: Image insert error slide {i} (query: {slide_data.get('image_search_query')}): {e}")
                traceback.print_exc()
                # Check first few bytes to debug 'cannot identify image file' errors
                try:
                    stream.seek(0)
                    header = stream.read(20)
                    print(f"--- DEBUG: First 20 bytes of failed image: {header}")
                except Exception:
                    pass

        # --- Chart or Table ---
        if slide_data.get("chart"):
            _add_chart(slide, slide_data["chart"], visual_x, visual_y, visual_w, visual_h)
        elif slide_data.get("table"):
            _add_table(slide, slide_data["table"], visual_x, visual_y, visual_w, visual_h)

        # Move off-screen any unused layout placeholders (prevents auto-resize)
        if use_image:
            for ph in slide.shapes.placeholders:
                if ph.placeholder_format.idx not in (0, 1):
                    ph.left = Inches(-20)

    # Save
    output_path = os.path.abspath(filename)
    prs.save(output_path)
    return output_path


# ===========================================================================
# IN-PLACE PPT PATCHING ENGINE
# Preserves original design; only modifies text/visuals as instructed.
# ===========================================================================

from pptx.enum.shapes import MSO_SHAPE_TYPE

# Shape types that count as "visuals" for layout decisions
_VISUAL_MSO_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}


def _count_visuals(slide):
    """Return (pictures, charts, tables) counts for a slide."""
    pics   = sum(1 for s in slide.shapes if s.shape_type in _VISUAL_MSO_TYPES)
    charts = sum(1 for s in slide.shapes if s.has_chart)
    tables = sum(1 for s in slide.shapes if s.has_table)
    return pics, charts, tables


def _get_body_shape(slide):
    """Find the largest non-title text frame shape on the slide."""
    title = slide.shapes.title
    candidates = [s for s in slide.shapes if s.has_text_frame and s != title]
    if not candidates:
        return None
    return max(candidates, key=lambda s: s.width * s.height)


def _sample_font_props(text_frame):
    """Sample font properties from the first available run in a text frame."""
    props = {}
    for para in text_frame.paragraphs:
        if para.runs:
            run = para.runs[0]
            try:
                props['size'] = run.font.size
            except Exception:
                pass
            try:
                props['bold'] = run.font.bold
            except Exception:
                pass
            try:
                props['italic'] = run.font.italic
            except Exception:
                pass
            try:
                props['name'] = run.font.name
            except Exception:
                pass
            try:
                if run.font.color and run.font.color.type:
                    props['color'] = run.font.color.rgb
            except Exception:
                pass
            break
    return props


def _sample_and_rewrite_text(text_frame, new_items):
    """
    Update a text frame's content while preserving the original font style.
    Only the paragraph text is changed — the shape, fill, position are untouched.
    """
    props = _sample_font_props(text_frame)
    text_frame.clear()  # Removes paragraph content only, NOT the shape itself

    for idx, item in enumerate(new_items):
        text  = item.get("text", "")
        level = item.get("level", 0)

        para = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        para.text  = text
        para.level = level

        # Re-apply sampled font style
        if props.get('size'):
            para.font.size = props['size']
        if props.get('bold') is not None:
            para.font.bold = props['bold']
        if props.get('italic') is not None:
            para.font.italic = props['italic']
        if props.get('name'):
            para.font.name = props['name']
        if props.get('color'):
            para.font.color.rgb = props['color']


def _remove_nth_shape_of_type(slide, shape_type_str, index=0):
    """
    Remove the Nth shape of the specified type from a slide via XML manipulation.
    shape_type_str: 'picture', 'chart', 'table', 'smartart'
    Returns True if found and removed, False otherwise.
    """
    matching = []
    for shape in slide.shapes:
        if shape_type_str == 'picture' and shape.shape_type in _VISUAL_MSO_TYPES:
            matching.append(shape)
        elif shape_type_str == 'chart' and shape.has_chart:
            matching.append(shape)
        elif shape_type_str == 'table' and shape.has_table:
            matching.append(shape)
        elif shape_type_str == 'smartart' and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            matching.append(shape)

    if index < len(matching):
        sp_el = matching[index]._element
        sp_el.getparent().remove(sp_el)
        print(f"--- PATCH: Removed {shape_type_str}[{index}]")
        return True

    print(f"--- PATCH: {shape_type_str}[{index}] not found (only {len(matching)} available)")
    return False


def _reflow_slide_content(slide, prs, was_visual_added=False, newly_added_shape=None):
    """
    Auto-resize the text body after a visual element is added or removed:
      - If no visuals remain on slide → expand text body to full available width
      - If a visual was just added AND text was previously full-width → shrink text
        to left 50% and position the new visual in the right 50%
    """
    W            = prs.slide_width
    RIGHT_MARGIN = Inches(0.5)
    GAP          = Inches(0.25)

    pics, charts, tables = _count_visuals(slide)
    total_visuals = pics + charts + tables

    body_shape = _get_body_shape(slide)
    if not body_shape:
        return

    available_width = W - body_shape.left - RIGHT_MARGIN
    half_width      = int(available_width * 0.50)

    if total_visuals == 0:
        # All visuals removed — expand text to full available width
        body_shape.width = int(available_width)
        print(f"--- REFLOW: No visuals remaining, expanded text to full width")

    elif was_visual_added and newly_added_shape is not None:
        # A visual was just added — check if text needs to be shrunk
        is_currently_full_width = body_shape.width > int(available_width * 0.60)
        if is_currently_full_width:
            body_shape.width = half_width
            print(f"--- REFLOW: Shrunk text to left 50% to make room for new visual")

        # Position the newly added shape in the right half
        visual_left   = body_shape.left + body_shape.width + int(GAP)
        visual_top    = body_shape.top
        visual_width  = W - visual_left - RIGHT_MARGIN
        visual_height = body_shape.height

        try:
            newly_added_shape.left   = int(visual_left)
            newly_added_shape.top    = int(visual_top)
            newly_added_shape.width  = int(visual_width)
            newly_added_shape.height = int(visual_height)
            print(f"--- REFLOW: Positioned new visual in right 50%")
        except Exception as e:
            print(f"--- REFLOW: Could not reposition new shape: {e}")


# ---------------------------------------------------------------------------
# extract_slide_inventory — read text + shape counts from an uploaded PPTX
# ---------------------------------------------------------------------------

def extract_slide_inventory(pptx_path):
    """
    Extract text content and visual element counts from each slide.
    Returns a list of slide dicts with:
      - heading: str
      - content: [{text, level}]
      - shapes: {images, charts, tables, smartarts}
    """
    prs    = Presentation(pptx_path)
    slides = []

    for slide in prs.slides:
        # Title
        heading = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            heading = slide.shapes.title.text_frame.text.strip()
        if not heading:
            heading = "Untitled Slide"

        # Body text
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        content.append({"text": text, "level": para.level})

        # Shape counts
        pics      = sum(1 for s in slide.shapes if s.shape_type in _VISUAL_MSO_TYPES)
        charts    = sum(1 for s in slide.shapes if s.has_chart)
        tables    = sum(1 for s in slide.shapes if s.has_table)
        smartarts = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP)

        slides.append({
            "heading": heading,
            "content": content,
            "shapes": {
                "images":    pics,
                "charts":    charts,
                "tables":    tables,
                "smartarts": smartarts
            }
        })

    return slides


# ---------------------------------------------------------------------------
# patch_ppt — main driver: apply structured action list to an existing PPTX
# ---------------------------------------------------------------------------

def patch_ppt(original_path, slide_actions, output_filename):
    """
    Patches an existing PPTX file based on a per-slide action list.
    Everything not referenced by an action is completely preserved.

    slide_actions: list of {
        "slide_index": int (0-based),
        "actions": [
            {"type": "update_text", "heading": str, "content": [{text, level}]},
            {"type": "remove_image",   "index": int},
            {"type": "remove_chart",   "index": int},
            {"type": "remove_table",   "index": int},
            {"type": "remove_smartart","index": int},
            {"type": "add_image",  "query": str},
            {"type": "add_chart",  "chart": {type, title, categories, series}},
            {"type": "add_table",  "table": [[...],[...]]},
        ]
    }
    """
    print(f"--- PATCH: Loading original from {original_path}")
    prs = Presentation(original_path)
    W   = prs.slide_width
    H   = prs.slide_height

    for slide_action in slide_actions:
        slide_idx = slide_action.get("slide_index", 0)
        actions   = slide_action.get("actions", [])

        if slide_idx >= len(prs.slides):
            print(f"--- PATCH: slide_index {slide_idx} out of range ({len(prs.slides)} slides), skipping")
            continue

        slide        = prs.slides[slide_idx]
        title_shape  = slide.shapes.title
        body_shape   = _get_body_shape(slide)

        newly_added_visual = None
        was_visual_added   = False

        for action in actions:
            atype = action.get("type", "")
            print(f"--- PATCH: slide {slide_idx} → {atype}")

            # ── Update text ──────────────────────────────────────────────────
            if atype == "update_text":
                new_heading = action.get("heading")
                new_content = action.get("content", [])

                if new_heading and title_shape and title_shape.has_text_frame:
                    _sample_and_rewrite_text(
                        title_shape.text_frame,
                        [{"text": new_heading, "level": 0}]
                    )

                # Re-fetch body_shape in case slide was modified
                body_shape = _get_body_shape(slide)
                if new_content and body_shape:
                    _sample_and_rewrite_text(body_shape.text_frame, new_content)

            # ── Remove shapes ────────────────────────────────────────────────
            elif atype == "remove_image":
                _remove_nth_shape_of_type(slide, "picture", action.get("index", 0))

            elif atype == "remove_chart":
                _remove_nth_shape_of_type(slide, "chart", action.get("index", 0))

            elif atype == "remove_table":
                _remove_nth_shape_of_type(slide, "table", action.get("index", 0))

            elif atype == "remove_smartart":
                _remove_nth_shape_of_type(slide, "smartart", action.get("index", 0))

            # ── Add image ────────────────────────────────────────────────────
            elif atype == "add_image":
                query = action.get("query", "")
                if query:
                    img_stream = fetch_image(query)
                    if img_stream:
                        try:
                            # Placeholder position — reflow will correct this
                            temp_shape = slide.shapes.add_picture(
                                img_stream,
                                Inches(7), Inches(1.5), Inches(5), Inches(4)
                            )
                            newly_added_visual = temp_shape
                            was_visual_added   = True
                            print(f"--- PATCH: Added image for query '{query}'")
                        except Exception as e:
                            print(f"--- PATCH: Error adding image: {e}")
                    else:
                        print(f"--- PATCH: No image found for query '{query}'")

            # ── Add chart ────────────────────────────────────────────────────
            elif atype == "add_chart":
                chart_data = action.get("chart", {})
                if chart_data:
                    # Placeholder position — reflow will correct this
                    chart_shape = _add_chart(
                        slide, chart_data,
                        Inches(7), Inches(1.5), Inches(5.5), Inches(4.5)
                    )
                    if chart_shape:
                        newly_added_visual = chart_shape
                        was_visual_added   = True

            # ── Add table ────────────────────────────────────────────────────
            elif atype == "add_table":
                table_data = action.get("table", [])
                if table_data:
                    table_shape = _add_table(
                        slide, table_data,
                        Inches(7), Inches(1.5), Inches(5.5), Inches(4.5)
                    )
                    if table_shape:
                        newly_added_visual = table_shape
                        was_visual_added   = True

            else:
                print(f"--- PATCH: Unknown action type '{atype}', skipping")

        # ── Auto-reflow after all actions on this slide ───────────────────
        _reflow_slide_content(slide, prs, was_visual_added, newly_added_visual)

    output_path = os.path.abspath(output_filename)
    prs.save(output_path)
    print(f"--- PATCH: Saved patched file → {output_path}")
    return output_path
