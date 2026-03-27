import ppt_utils
import os

def test_layout():
    print("Testing PPT layout logic...")
    
    mock_data = {
        "title": "Layout Test Presentation",
        "theme": {"font": "Calibri"},
        "slides": [
            {
                "heading": "Slide 1: Text Only (Full Width)",
                "content": [
                    {"text": "This slide should use the full width of the layout because it has no visual elements.", "level": 0},
                    {"text": "Notice how the title and body span the entire slide.", "level": 0}
                ]
            },
            {
                "heading": "Slide 2: Text and Chart (Two Column)",
                "content": [
                    {"text": "This slide should use a two-column layout because it contains a chart.", "level": 0},
                    {"text": "The title uses 65% width to allow more room for long headings.", "level": 1}
                ],
                "chart": {
                    "type": "column",
                    "title": "Sample Chart",
                    "categories": ["A", "B"],
                    "series": [{"name": "Series 1", "values": [10, 20]}]
                }
            },
            {
                "heading": "Slide 3: This is a Very Long Title That Should Not Wrap Prematurely Even in Two Column Mode",
                "content": [
                    {"text": "Even with an image/table, the title now has 65% of the slide width.", "level": 0}
                ],
                "table": [
                    ["Feature", "Status"],
                    ["Dynamic Layout", "Active"],
                    ["Expanded Titles", "Active"]
                ]
            }
        ],
        "design_prefs": {
            "template": "professional_blue.pptx",
            "title_font_size": 32,
            "body_font_size": 20
        }
    }
    
    filename = "verify_layout_test.pptx"
    try:
        output_path = ppt_utils.create_ppt(mock_data, filename=filename, image_mode="none")
        print(f"\n--- VERIFICATION SUCCESS ---")
        print(f"PPT saved to: {output_path}")
        
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"Total Slides: {len(prs.slides)}")
        print("\nSlide Headings:")
        for i, slide in enumerate(prs.slides):
            print(f"{i+1}. {slide.shapes.title.text}")
        
        print(f"\nACTION: Please open '{filename}' to visually confirm that:")
        print("1. Slide 2 (Text Only) uses the FULL width of the slide.")
        print("2. Slide 3 & 4 (Chart/Table) use the 2-column layout.")
        print("3. The title in Slide 4 does not wrap prematurely.")
            
    except Exception as e:
        print(f"Error during layout test: {e}")

if __name__ == "__main__":
    test_layout()
